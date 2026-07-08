from engine.backend_pptx import new_presentation, add_slide
from engine.layout import text, rect, chip, Slide
from engine.theme import load_theme
from pptx.util import Emu

T = load_theme(None)

def test_presentation_is_16_9():
    prs = new_presentation(T)
    assert round(Emu(prs.slide_width).inches, 2) == 13.33
    assert round(Emu(prs.slide_height).inches, 2) == 7.5

def test_add_slide_creates_shapes_with_text():
    prs = new_presentation(T)
    s = Slide("title", [
        rect(0, 0, 1, 1, fill="#141A2E"),
        text(0.1, 0.4, 0.6, 0.1, "КВАРТАЛ GDS", color="#FFFFFF", size_pt=25, bold=True),
        chip(0.1, 0.6, 0.12, 0.05, "BE·Go", fill="#2E4B7A"),
    ], footer="GDS · Q3", page=1)
    add_slide(prs, s, T)
    assert len(prs.slides) == 1
    texts = [sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame]
    assert any("КВАРТАЛ GDS" in t for t in texts)
    assert any("BE·Go" in t for t in texts)
    assert any("GDS · Q3" in t for t in texts)  # footer drawn

def test_saves_to_disk(tmp_path):
    prs = new_presentation(T)
    add_slide(prs, Slide("t", [text(0, 0, 1, 0.1, "X")], page=1), T)
    out = tmp_path / "d.pptx"
    prs.save(out)
    assert out.exists() and out.stat().st_size > 0
