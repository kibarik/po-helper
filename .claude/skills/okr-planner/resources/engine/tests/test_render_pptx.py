from pathlib import Path
from engine.render_pptx import render_pptx_file
from pptx import Presentation

FIX = Path(__file__).parent / "fixtures" / "minimal_deck.yaml"

def test_render_pptx_file(tmp_path):
    out = render_pptx_file(FIX, tmp_path / "deck.pptx")
    assert out.exists()
    prs = Presentation(str(out))
    assert len(prs.slides) >= 6   # title + big_picture + glossary + divider + lifecycle + detail + closings
    all_text = " ".join(sh.text_frame.text for sl in prs.slides
                        for sh in sl.shapes if sh.has_text_frame)
    assert "Продажа мероприятий TicketsCloud" in all_text
    assert "БАЗИС" in all_text
