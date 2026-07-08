from __future__ import annotations
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from .layout import Element, Slide
from .theme import Theme

_EMU_W = 12192000
_EMU_H = 6858000
_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}


def _rgb(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr.lstrip("#").upper())


def new_presentation(theme: Theme) -> Presentation:
    prs = Presentation()
    prs.slide_width = Emu(_EMU_W)
    prs.slide_height = Emu(_EMU_H)
    return prs


def _emu(e: Element):
    return (Emu(int(e.fx * _EMU_W)), Emu(int(e.fy * _EMU_H)),
            Emu(int(e.fw * _EMU_W)), Emu(int(e.fh * _EMU_H)))


def _fill_shape(shape, hexstr):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hexstr)
    shape.line.fill.background()


def _write_text(tf, e: Element, theme: Theme):
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(45720)
    tf.margin_top = tf.margin_bottom = Emu(18288)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if e.valign == "middle" else MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = _ALIGN.get(e.align, PP_ALIGN.LEFT)
    runs = e.runs or [(e.text, e.color or theme.body, e.bold)]
    for txt, color, bold in runs:
        r = p.add_run()
        r.text = txt
        r.font.name = e.font or theme.body_font
        r.font.size = Pt(e.size_pt or 12)
        r.font.bold = bold
        r.font.color.rgb = _rgb(color or theme.body)


def add_slide(prs: Presentation, s: Slide, theme: Theme) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    for e in s.elements:
        x, y, w, h = _emu(e)
        if e.kind == "rect":
            shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if e.radius
                                         else MSO_SHAPE.RECTANGLE, x, y, w, h)
            _fill_shape(shp, e.fill or "#FFFFFF")
            if e.text or e.runs:
                _write_text(shp.text_frame, e, theme)
        elif e.kind == "chip":
            shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            _fill_shape(shp, e.fill or theme.card_bg)
            e2 = Element("text", e.fx, e.fy, e.fw, e.fh, text=e.text, runs=e.runs,
                         color=e.color or "#FFFFFF", font=e.font, size_pt=e.size_pt or 9,
                         bold=True, align="center", valign="middle")
            _write_text(shp.text_frame, e2, theme)
        else:  # text
            box = slide.shapes.add_textbox(x, y, w, h)
            _write_text(box.text_frame, e, theme)
    if s.footer or s.page:
        fl = slide.shapes.add_textbox(Emu(int(0.042 * _EMU_W)), Emu(int(0.93 * _EMU_H)),
                                      Emu(int(0.6 * _EMU_W)), Emu(int(0.05 * _EMU_H)))
        _write_text(fl.text_frame, Element("text", 0, 0, 0, 0, text=s.footer,
                                           color=theme.muted, size_pt=9.5), theme)
        fr = slide.shapes.add_textbox(Emu(int(0.9 * _EMU_W)), Emu(int(0.93 * _EMU_H)),
                                      Emu(int(0.06 * _EMU_W)), Emu(int(0.05 * _EMU_H)))
        _write_text(fr.text_frame, Element("text", 0, 0, 0, 0, text=str(s.page),
                                           color=theme.muted, size_pt=9.5, align="right"), theme)
